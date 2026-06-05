#!/usr/bin/env python3
"""
PPT Exam Review — Content Extraction Engine
=============================================
Extracts structured Markdown from PPTX and PDF files for exam review processing.

Supports:
  - PPTX: full text hierarchy, tables, speaker notes, bold/heading detection
  - PDF:  text-based PDFs via pdfplumber (tables + text)
  - OCR:  scanned/image-based PDFs via pytesseract (optional)

Usage:
  python extract_ppt.py <filepath> [options]

Output:
  - Markdown (default): ready for LLM consumption, one file per slide
  - JSON (--format json): structured data for programmatic use

Dependencies:
  pip install python-pptx pdfplumber
  # OCR optional: pip install pytesseract Pillow && install tesseract-ocr
"""

import argparse
import json
import os
import re
import sys
from pathlib import Path

# Fix Windows GBK console output
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')


# ── Utilities ──────────────────────────────────────────────────────────────

def coalesce_lines(texts: list[dict]) -> list[dict]:
    """Merge consecutive same-level text items into paragraphs."""
    if not texts:
        return []
    merged = [texts[0]]
    for item in texts[1:]:
        if item["level"] == merged[-1]["level"] and item["level"] > 0 and not item["bold"]:
            merged[-1]["text"] += " " + item["text"]
        else:
            merged.append(item)
    return merged


def detect_heading_level(text: str, font_size_pt: float | None, is_bold: bool) -> int:
    """Heuristic heading detection.

    Level 0 = body text (default)
    Level 1 = slide title or large heading (>=24pt or short bold with uppercase start)
    Level 2 = section heading (>=18pt or bold)
    Level 3 = subsection
    """
    if font_size_pt is not None:
        if font_size_pt >= 28:
            return 1
        if font_size_pt >= 20:
            return 2
        if font_size_pt >= 16:
            return 3
    if is_bold and len(text) < 60:
        if any(kw in text for kw in ("第", "章", "节", "Part", "Chapter", "Section", "§")):
            return 1
        return 2
    return 0


def clean_text(text: str) -> str:
    """Normalize whitespace and remove control chars."""
    text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f]', '', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


def clean_cid(text: str) -> str:
    """Remove (cid:NNN) artifacts from formula extraction."""
    text = re.sub(r'\(cid:\d+\)', '', text)
    return text.strip()


def is_garbled(text: str, threshold: float = 0.15) -> bool:
    """Heuristic: if >threshold of chars are encoding artifacts, text is garbled."""
    if not text or len(text) < 10:
        return False
    artifact_chars = sum(1 for c in text if ord(c) > 127 and c in '�□')
    cid_hits = len(re.findall(r'\(cid:\d+\)', text))
    return (artifact_chars / len(text) > threshold) or (cid_hits > 5)


# ── PPTX Extraction ───────────────────────────────────────────────────────

def extract_pptx(filepath: str) -> list[dict]:
    """Extract all slide content from a .pptx file."""
    from pptx import Presentation
    from pptx.util import Pt, Emu

    prs = Presentation(filepath)
    slides_data = []

    for idx, slide in enumerate(prs.slides, 1):
        slide_info = {
            "slide_number": idx,
            "texts": [],
            "tables": [],
            "notes": "",
            "shapes_count": len(slide.shapes),
        }

        # Speaker notes (high priority — often contains exam hints)
        try:
            if slide.has_notes_slide:
                notes_frame = slide.notes_slide.notes_text_frame
                raw = notes_frame.text.strip()
                if raw:
                    slide_info["notes"] = clean_text(raw)
        except Exception:
            pass

        for shape in slide.shapes:
            # ── Text frames ──
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = clean_text(para.text)
                    if not text:
                        continue
                    font_size = None
                    is_bold = False
                    try:
                        if para.font and para.font.size:
                            font_size = para.font.size.pt
                        if para.font and para.font.bold:
                            is_bold = True
                    except Exception:
                        pass
                    level = detect_heading_level(text, font_size, is_bold)
                    # Override: if this is the first shape on slide, it's likely the title
                    if len(slide_info["texts"]) == 0 and shape == slide.shapes[0]:
                        level = min(level, 1)  # at most level 1
                        if is_bold:
                            level = 1
                    slide_info["texts"].append({
                        "text": text,
                        "level": level,
                        "bold": is_bold,
                    })

            # ── Tables ──
            if shape.has_table:
                table = shape.table
                table_data = []
                for row_idx, row in enumerate(table.rows):
                    cells = []
                    for cell in row.cells:
                        cell_text = clean_text(cell.text)
                        cells.append(cell_text)
                    table_data.append(cells)
                if table_data:
                    slide_info["tables"].append(table_data)

            # ── Group shapes (recursive) ──
            if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
                try:
                    for child in shape.shapes:
                        if child.has_text_frame:
                            for para in child.text_frame.paragraphs:
                                text = clean_text(para.text)
                                if text:
                                    slide_info["texts"].append({
                                        "text": text,
                                        "level": 2,
                                        "bold": False,
                                    })
                except Exception:
                    pass

        slide_info["texts"] = coalesce_lines(slide_info["texts"])
        slides_data.append(slide_info)

    return slides_data


# ── PDF Extraction ─────────────────────────────────────────────────────────

def extract_pdf_pymupdf(filepath: str) -> list[dict]:
    """Extract PDF content using PyMuPDF (better for CJK with CID/Identity-H encodings)."""
    try:
        import fitz
    except ImportError:
        return None  # signal to use fallback

    doc = fitz.open(filepath)
    pages_data = []

    for idx, page_num in enumerate(range(len(doc)), 1):
        page = doc[page_num]
        page_info = {
            "slide_number": idx,
            "texts": [],
            "tables": [],
            "notes": "",
            "shapes_count": 0,
        }

        # Extract text blocks with position info
        blocks = page.get_text("dict", flags=fitz.TEXT_PRESERVE_WHITESPACE)["blocks"]
        for block in blocks:
            if block["type"] != 0:  # skip images
                continue
            for line in block["lines"]:
                text_parts = []
                max_size = 0
                is_bold = False
                for span in line["spans"]:
                    text_parts.append(span["text"])
                    max_size = max(max_size, span["size"])
                    if span["flags"] & 2:  # bold flag
                        is_bold = True

                line_text = clean_text("".join(text_parts))
                if not line_text:
                    continue

                level = detect_heading_level(line_text, max_size, is_bold)
                page_info["texts"].append({
                    "text": line_text,
                    "level": level,
                    "bold": is_bold,
                })

        page_info["texts"] = coalesce_lines(page_info["texts"])
        pages_data.append(page_info)

    doc.close()
    return pages_data


def extract_pdf(filepath: str) -> list[dict]:
    """Extract content from a text-based PDF file.

    Uses pdfplumber (primary) with PyMuPDF fallback for CID-encoded CJK PDFs.
    """
    try:
        import pdfplumber
    except ImportError:
        print("ERROR: PDF extraction requires pdfplumber.", file=sys.stderr)
        print("  pip install pdfplumber", file=sys.stderr)
        sys.exit(1)

    pages_data = []
    garbled_detected = False

    with pdfplumber.open(filepath) as pdf:
        for idx, page in enumerate(pdf.pages, 1):
            page_info = {
                "slide_number": idx,
                "texts": [],
                "tables": [],
                "notes": "",
                "shapes_count": 0,
            }

            # Extract text
            raw_text = page.extract_text() or ""
            for line in raw_text.split("\n"):
                line = clean_text(line)
                if not line:
                    continue
                # Heuristic: short all-caps / large font-like lines → heading
                is_heading = bool(
                    (line.isupper() and len(line) > 3 and len(line) < 80)
                    or re.match(r'^第[一二三四五六七八九十\d]+[章节部]', line)
                    or re.match(r'^(Chapter|Section|Part)\s+\d', line, re.I)
                )
                level = 1 if is_heading else 0
                page_info["texts"].append({
                    "text": line,
                    "level": level,
                    "bold": is_heading,
                })

            # Check for garbled content
            if not garbled_detected and is_garbled(raw_text):
                garbled_detected = True

            # Extract tables
            raw_tables = page.extract_tables()
            if raw_tables:
                for table in raw_tables:
                    clean_table = []
                    for row in table:
                        clean_row = [clean_text(c) if c else "" for c in row]
                        if any(cell for cell in clean_row):
                            clean_table.append(clean_row)
                    if clean_table:
                        page_info["tables"].append(clean_table)

            page_info["texts"] = coalesce_lines(page_info["texts"])
            pages_data.append(page_info)

    # If garbled content detected, fall back to PyMuPDF
    if garbled_detected:
        print("  [WARN] pdfplumber produced garbled text; retrying with PyMuPDF...", file=sys.stderr)
        fallback = extract_pdf_pymupdf(filepath)
        if fallback:
            return fallback

    return pages_data


# ── OCR Extraction (scanned PDFs/images) ──────────────────────────────────

def extract_ocr(filepath: str) -> list[dict]:
    """Extract text from scanned/image-based documents via OCR.

    Requires: pip install pytesseract Pillow
    Also requires tesseract-ocr installed on the system.
    """
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        print("ERROR: OCR requires pytesseract and Pillow.", file=sys.stderr)
        print("  pip install pytesseract Pillow", file=sys.stderr)
        print("  Also install tesseract-ocr: https://github.com/tesseract-ocr/tesseract", file=sys.stderr)
        sys.exit(1)

    ext = Path(filepath).suffix.lower()

    # Single image
    if ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        img = Image.open(filepath)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        lines = [clean_text(l) for l in text.split("\n") if clean_text(l)]
        return [{
            "slide_number": 1,
            "texts": [{"text": l, "level": 0, "bold": False} for l in lines],
            "tables": [],
            "notes": "[OCR extracted — verify accuracy]",
            "shapes_count": 0,
        }]

    # PDF → convert pages to images → OCR each
    try:
        from pdf2image import convert_from_path
    except ImportError:
        print("ERROR: PDF OCR requires pdf2image.", file=sys.stderr)
        print("  pip install pdf2image", file=sys.stderr)
        sys.exit(1)

    print("  Converting PDF pages to images for OCR...", file=sys.stderr)
    images = convert_from_path(filepath, dpi=300)
    pages_data = []
    for idx, img in enumerate(images, 1):
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        lines = [clean_text(l) for l in text.split("\n") if clean_text(l)]
        pages_data.append({
            "slide_number": idx,
            "texts": [{"text": l, "level": 0, "bold": False} for l in lines],
            "tables": [],
            "notes": "[OCR extracted — verify accuracy]",
            "shapes_count": 0,
        })

    return pages_data


# ── Output formatters ─────────────────────────────────────────────────────

def to_markdown(slides_data: list[dict], include_notes: bool = True) -> str:
    """Convert extracted slide data to structured Markdown."""
    empty_count = sum(1 for s in slides_data if not s["texts"] and not s["tables"])
    lines = []
    for slide in slides_data:
        lines.append(f"## Slide {slide['slide_number']}")
        lines.append("")

        if not slide["texts"] and not slide["tables"]:
            lines.append("_[No extractable text content on this slide]_")
            lines.append("")
            continue

        for item in slide["texts"]:
            prefix = "#" * min(item["level"] + 2, 6)  # start at ## for slide-level
            text = clean_cid(item["text"])  # clean (cid:xxx) artifacts
            if not text:
                continue
            if item["bold"] and item["level"] <= 1:
                lines.append(f"{prefix} **{text}**")
            elif item["level"] >= 2:
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)
            lines.append("")

        for table in slide["tables"]:
            if not table:
                continue
            header = "| " + " | ".join(table[0]) + " |"
            sep = "| " + " | ".join(["---"] * len(table[0])) + " |"
            lines.append(header)
            lines.append(sep)
            for row in table[1:]:
                lines.append("| " + " | ".join(row) + " |")
            lines.append("")

        if include_notes and slide.get("notes"):
            lines.append("> **📝 Speaker Notes:** " + slide["notes"])
            lines.append("")

    return "\n".join(lines)


def to_json(slides_data: list[dict]) -> str:
    """Output as pretty-printed JSON."""
    return json.dumps(slides_data, ensure_ascii=False, indent=2)


# ── Summary ────────────────────────────────────────────────────────────────

def print_summary(slides_data: list[dict]) -> None:
    """Print a quick overview of extracted content."""
    total_slides = len(slides_data)
    total_texts = sum(len(s["texts"]) for s in slides_data)
    total_tables = sum(len(s["tables"]) for s in slides_data)
    slides_with_notes = sum(1 for s in slides_data if s.get("notes"))

    print(f"\n{'='*50}", file=sys.stderr)
    print(f"📊  Extraction Summary", file=sys.stderr)
    print(f"{'='*50}", file=sys.stderr)
    print(f"  Slides:        {total_slides}", file=sys.stderr)
    print(f"  Text blocks:   {total_texts}", file=sys.stderr)
    print(f"  Tables:        {total_tables}", file=sys.stderr)
    print(f"  With notes:    {slides_with_notes}", file=sys.stderr)
    print(f"{'='*50}\n", file=sys.stderr)


# ── CLI ────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(
        description="Extract structured content from PPT/PDF for exam review processing.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=(
            "Examples:\n"
            "  python extract_ppt.py lecture.pptx\n"
            "  python extract_ppt.py lecture.pdf --format json\n"
            "  python extract_ppt.py scanned.pdf --ocr\n"
            "  python extract_ppt.py slide.png --ocr\n"
        ),
    )
    parser.add_argument("filepath", help="Path to PPTX, PDF, or image file")
    parser.add_argument("--format", choices=["markdown", "json"], default="markdown",
                       help="Output format (default: markdown)")
    parser.add_argument("--no-notes", action="store_true",
                       help="Exclude speaker notes from output")
    parser.add_argument("--ocr", action="store_true",
                       help="Force OCR mode (for scanned/image-based content)")
    parser.add_argument("--quiet", "-q", action="store_true",
                       help="Suppress summary output")

    args = parser.parse_args()

    filepath = Path(args.filepath)
    if not filepath.exists():
        print(f"Error: File not found: {filepath}", file=sys.stderr)
        sys.exit(1)

    ext = filepath.suffix.lower()

    # Route to appropriate extractor
    if args.ocr:
        slides_data = extract_ocr(str(filepath))
    elif ext == ".pptx":
        slides_data = extract_pptx(str(filepath))
    elif ext == ".pdf":
        slides_data = extract_pdf(str(filepath))
    elif ext in (".png", ".jpg", ".jpeg", ".tiff", ".bmp"):
        slides_data = extract_ocr(str(filepath))
    else:
        print(f"Error: Unsupported format '{ext}'. Supported: .pptx, .pdf, .png, .jpg", file=sys.stderr)
        sys.exit(1)

    # Check for empty results (possible scan/image-only PDF)
    if not args.ocr and slides_data:
        empty_slides = sum(1 for s in slides_data if not s["texts"] and not s["tables"])
        if empty_slides == len(slides_data):
            print("  [HINT] All slides appear empty. This file may be image-based/scanned.", file=sys.stderr)
            print("  [HINT] Try: pip install pytesseract pdf2image Pillow && extract_ppt.py --ocr", file=sys.stderr)

    if not args.quiet:
        print_summary(slides_data)

    if args.format == "json":
        print(to_json(slides_data))
    else:
        print(to_markdown(slides_data, include_notes=not args.no_notes))


if __name__ == "__main__":
    main()
